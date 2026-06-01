"""
18650 Cylindrical Li-ion Battery Two-Node Thermal Model & Thermal Equivalent Circuit
=====================================================================================
A fresh, standalone script generating a publication-quality figure.

Run:
    pip install matplotlib numpy python-pptx
    python plot_battery_thermal.py

Output:
    cylindrical_battery_two_node_thermal_model.png  (300 dpi)
    cylindrical_battery_two_node_thermal_model.pptx (16:9 slide)
"""

import numpy as np
import matplotlib.pyplot as plt
from matplotlib.patches import Ellipse, FancyBboxPatch, FancyArrowPatch
from matplotlib.transforms import Affine2D

# ============================================================
# CONFIGURATION
# ============================================================
CFG = {
    "dpi": 300,
    "fig_w": 14.5,
    "fig_h": 7.0,
    "font": "serif",
    "mathtext": "cm",
    # battery geometry (data coords)
    "bat_cx": 0.0,
    "bat_cy": 0.0,
    "bat_w": 1.1,
    "bat_h": 2.6,
    "bat_tilt": 0.17,
    # colors
    "shell_pink": (1.0, 0.68, 0.76),
    "core_pink": (1.0, 0.85, 0.90),
    "metal_base": (0.72, 0.72, 0.76),
    "metal_bright": (0.94, 0.94, 0.97),
    "insul_black": (0.10, 0.10, 0.10),
    "tc_color": (0.50, 0.0, 0.50),
    "ts_color": (0.85, 0.10, 0.10),
    "rc_color": (1.0, 0.78, 0.85),
    "rs_color": (1.0, 0.55, 0.15),
    "wire": "black",
    "arrow_blue": (0.35, 0.60, 0.92, 0.65),
}

plt.rcParams.update({
    "font.family": CFG["font"],
    "mathtext.fontset": CFG["mathtext"],
    "font.size": 10,
})


# ============================================================
# DRAWING PRIMITIVES
# ============================================================

def _gradient_rect_v(ax, x, y, w, h, color_func, n=30, zorder=1):
    """Fill a rectangle with vertical gradient strips using color_func(frac)."""
    for i in range(n):
        frac = (i + 0.5) / n
        xi = x + i * w / n
        c = color_func(frac)
        ax.add_patch(plt.Rectangle((xi, y), w / n, h, fc=c, ec="none", zorder=zorder))


def draw_battery(ax):
    """Draw the 18650 cylindrical battery with 3D appearance."""
    cx, cy = CFG["bat_cx"], CFG["bat_cy"]
    W, H = CFG["bat_w"], CFG["bat_h"]
    a = W / 2
    b = a * CFG["bat_tilt"]
    top_y = cy + H / 2
    bot_y = cy - H / 2

    # -- cylinder body gradient --
    def shell_color(frac):
        d = abs(frac - 0.5) / 0.5
        r, g, bl = CFG["shell_pink"]
        alpha = 0.30 + 0.30 * (1 - d ** 2)
        return (r, g - 0.12 * d, bl - 0.05 * d, alpha)

    _gradient_rect_v(ax, cx - a, bot_y, W, H, shell_color, n=50, zorder=1)

    # -- core hint --
    ca = a * 0.55
    def core_color(frac):
        d = abs(frac - 0.5) / 0.5
        return (*CFG["core_pink"], 0.12 + 0.18 * (1 - d ** 2))
    _gradient_rect_v(ax, cx - ca, bot_y + H * 0.08, 2 * ca, H * 0.84, core_color, n=20, zorder=1.5)

    # -- outlines --
    edge_c = (0.55, 0.28, 0.38, 0.85)
    ax.plot([cx - a, cx - a], [bot_y, top_y], color=edge_c, lw=1.4, zorder=3)
    ax.plot([cx + a, cx + a], [bot_y, top_y], color=edge_c, lw=1.4, zorder=3)

    # -- bottom ellipse --
    ax.add_patch(Ellipse((cx, bot_y), W, 2 * b, fc=(*CFG["shell_pink"], 0.5),
                         ec=edge_c, lw=1.2, zorder=2))

    # -- top ellipse --
    ax.add_patch(Ellipse((cx, top_y), W, 2 * b, fc=(*CFG["core_pink"], 0.85),
                         ec=edge_c, lw=1.2, zorder=4))

    # -- positive terminal cap --
    _draw_cap(ax, cx, top_y, a, b)

    return cx, cy, a, b, top_y, bot_y


def _draw_cap(ax, cx, top_y, a, b):
    """Draw the 3D metallic positive cap sitting on battery top."""
    cap_h = 0.11
    ca = a * 0.90
    cb = b * 0.90
    base_y = top_y
    face_y = base_y + cap_h

    # side wall gradient
    def metal_side(frac):
        d = abs(frac - 0.5) / 0.5
        v = 0.82 - 0.28 * d ** 1.3
        return (v, v, v + 0.02, 0.97)
    _gradient_rect_v(ax, cx - ca, base_y, 2 * ca, cap_h, metal_side, n=35, zorder=4.5)

    # side edges
    ec = (0.32, 0.32, 0.35)
    ax.plot([cx - ca, cx - ca], [base_y, face_y], color=ec, lw=1.3, zorder=4.6)
    ax.plot([cx + ca, cx + ca], [base_y, face_y], color=ec, lw=1.3, zorder=4.6)

    # top face (outer ring)
    ax.add_patch(Ellipse((cx, face_y), 2 * ca, 2 * cb,
                         fc=CFG["metal_base"], ec=ec, lw=1.2, zorder=5))
    # highlight
    ax.add_patch(Ellipse((cx - ca * 0.22, face_y + cb * 0.18), ca * 0.65, cb * 0.5,
                         fc=CFG["metal_bright"], ec="none", alpha=0.55, zorder=5.05))

    # insulation ring
    ia, ib = ca * 0.70, cb * 0.70
    ax.add_patch(Ellipse((cx, face_y), 2 * ia, 2 * ib,
                         fc=CFG["insul_black"], ec=(0.0, 0.0, 0.0), lw=0.8, zorder=5.1))

    # central bump
    ba, bb = ca * 0.36, cb * 0.48
    ax.add_patch(Ellipse((cx, face_y + 0.012), 2 * ba * 1.08, 2 * bb * 1.08,
                         fc=(0.58, 0.58, 0.62), ec=(0.38, 0.38, 0.40), lw=0.6, zorder=5.2))
    ax.add_patch(Ellipse((cx, face_y + 0.022), 2 * ba, 2 * bb,
                         fc=(0.80, 0.80, 0.84), ec=(0.48, 0.48, 0.50), lw=0.7, zorder=5.3))
    ax.add_patch(Ellipse((cx - ba * 0.18, face_y + 0.032), ba * 0.55, bb * 0.45,
                         fc=(0.97, 0.97, 1.0), ec="none", alpha=0.7, zorder=5.4))

    # "+" label
    ax.text(cx, face_y + cb + 0.07, "+", fontsize=10, fontweight="bold",
            ha="center", va="center", zorder=6)


# ---- circuit element primitives ----

def draw_heat_source(ax, x, y, r=0.10, label=r"$Q_e$", lbl_dx=-0.22, lbl_dy=0):
    """Circle with sine wave."""
    ax.add_patch(plt.Circle((x, y), r, fc="white", ec="black", lw=1.6, zorder=10))
    t = np.linspace(-0.65 * r, 0.65 * r, 40)
    w = np.sin(t / r * 2.8 * np.pi) * r * 0.32
    ax.plot(t + x, w + y, "k-", lw=1.1, zorder=11)
    ax.text(x + lbl_dx, y + lbl_dy, label, fontsize=11, ha="center", va="center", zorder=11)


def draw_resistor_h(ax, x1, x2, y, color, label=""):
    """Horizontal resistor between (x1,y) and (x2,y)."""
    mx = (x1 + x2) / 2
    rl = (x2 - x1) * 0.48
    rh = 0.10
    ax.plot([x1, mx - rl / 2], [y, y], color=CFG["wire"], lw=1.8, zorder=9)
    ax.plot([mx + rl / 2, x2], [y, y], color=CFG["wire"], lw=1.8, zorder=9)
    ax.add_patch(FancyBboxPatch((mx - rl / 2, y - rh / 2), rl, rh,
                                boxstyle="round,pad=0.012", fc=color, ec="black", lw=1.3, zorder=10))
    if label:
        ax.text(mx, y + rh / 2 + 0.09, label, fontsize=11, ha="center", va="center", zorder=11)


def draw_resistor_v(ax, x, y1, y2, color, label="", lbl_side="right"):
    """Vertical resistor between (x,y1) top and (x,y2) bottom."""
    my = (y1 + y2) / 2
    rl = abs(y1 - y2) * 0.48
    rw = 0.10
    ax.plot([x, x], [y1, my + rl / 2], color=CFG["wire"], lw=1.8, zorder=9)
    ax.plot([x, x], [my - rl / 2, y2], color=CFG["wire"], lw=1.8, zorder=9)
    ax.add_patch(FancyBboxPatch((x - rw / 2, my - rl / 2), rw, rl,
                                boxstyle="round,pad=0.012", fc=color, ec="black", lw=1.3, zorder=10))
    if label:
        dx = 0.14 if lbl_side == "right" else -0.14
        ha = "left" if lbl_side == "right" else "right"
        ax.text(x + dx, my, label, fontsize=11, ha=ha, va="center", zorder=11)


def draw_capacitor_v(ax, x, y_top, y_bot, label="", lbl_side="right"):
    """Vertical capacitor."""
    my = (y_top + y_bot) / 2
    gap = 0.065
    pw = 0.13
    ax.plot([x, x], [y_top, my + gap], color=CFG["wire"], lw=1.8, zorder=9)
    ax.plot([x, x], [my - gap, y_bot], color=CFG["wire"], lw=1.8, zorder=9)
    ax.plot([x - pw, x + pw], [my + gap, my + gap], "k-", lw=2.6, zorder=10)
    ax.plot([x - pw, x + pw], [my - gap, my - gap], "k-", lw=2.6, zorder=10)
    if label:
        dx = 0.17 if lbl_side == "right" else -0.17
        ha = "left" if lbl_side == "right" else "right"
        ax.text(x + dx, my, label, fontsize=11, ha=ha, va="center", zorder=11)


def draw_ground_symbol(ax, x, y, size=0.09):
    """Standard ground symbol."""
    ax.plot([x, x], [y, y - size * 0.5], color=CFG["wire"], lw=1.8, zorder=9)
    for i, f in enumerate([1.0, 0.62, 0.28]):
        yy = y - size * 0.5 - i * size * 0.27
        ax.plot([x - size * f, x + size * f], [yy, yy], color=CFG["wire"], lw=1.8, zorder=9)


def draw_dot(ax, x, y, color, r=0.042, label="", pos="above"):
    """Node dot."""
    ax.add_patch(plt.Circle((x, y), r, fc=color, ec="black", lw=0.8, zorder=12))
    if label:
        offsets = {"above": (0, 0.14), "below": (0, -0.14), "left": (-0.14, 0), "right": (0.14, 0)}
        dx, dy = offsets.get(pos, (0, 0.14))
        ha = "center" if pos in ("above", "below") else ("right" if pos == "left" else "left")
        ax.text(x + dx, y + dy, label, fontsize=11, ha=ha, va="center", zorder=12)


# ============================================================
# LEFT PANEL: Battery + internal thermal network
# ============================================================

def build_left_panel(ax):
    ax.set_xlim(-1.3, 2.0)
    ax.set_ylim(-1.6, 2.0)
    ax.set_aspect("equal")
    ax.axis("off")

    # Draw battery
    cx, cy, a, b, top_y, bot_y = draw_battery(ax)

    # --- thermal network inside battery ---
    qx, qy = cx - 0.18, cy - 0.60
    tcx, tcy = cx - 0.18, cy + 0.30
    tsx, tsy = cx + 0.32, cy + 0.30
    gy = cy - 1.10  # ground y

    # heat source
    draw_heat_source(ax, qx, qy, r=0.095, lbl_dx=-0.24)

    # wires from Qe
    ax.plot([qx, qx], [qy + 0.095, tcy], color=CFG["wire"], lw=1.5, zorder=9)
    ax.plot([qx, qx], [qy - 0.095, gy], color=CFG["wire"], lw=1.5, zorder=9)

    # nodes
    draw_dot(ax, tcx, tcy, CFG["tc_color"], label=r"$T_c$", pos="above")
    draw_dot(ax, tsx, tsy, CFG["ts_color"], label=r"$T_s$", pos="above")

    # Rc horizontal between Tc and Ts
    draw_resistor_h(ax, tcx, tsx, tcy, CFG["rc_color"], label=r"$R_c$")

    # Cc below Tc
    draw_capacitor_v(ax, tcx, tcy - 0.10, gy + 0.05, label=r"$C_c$", lbl_side="left")

    # Rs below Ts
    draw_resistor_v(ax, tsx, tsy - 0.05, gy + 0.40, CFG["rs_color"], label=r"$R_s$", lbl_side="right")
    ax.plot([tsx, tsx], [gy + 0.40, gy], color=CFG["wire"], lw=1.4, zorder=9)

    # Cs parallel to Rs
    csx = tsx + 0.28
    ax.plot([tsx, csx], [tsy, tsy], color=CFG["wire"], lw=1.3, zorder=9)
    draw_capacitor_v(ax, csx, tsy - 0.10, gy + 0.05, label=r"$C_s$", lbl_side="right")
    ax.plot([csx, csx], [gy + 0.05, gy], color=CFG["wire"], lw=1.3, zorder=9)

    # ground rail
    ax.plot([qx - 0.08, csx + 0.08], [gy, gy], color=CFG["wire"], lw=1.6, zorder=9)
    gnd_mid = (tcx + tsx) / 2
    draw_ground_symbol(ax, gnd_mid, gy)
    ax.text(gnd_mid, gy - 0.20, r"$T_a$", fontsize=11, ha="center", va="top")

    # --- annotations ---
    ann_x = 1.15
    # Tc
    ax.plot([tcx + 0.06, ann_x], [tcy, tcy], ls="--", color=CFG["tc_color"], lw=1.1, zorder=8)
    ax.text(ann_x + 0.03, tcy + 0.09, "Internal Temperature", fontsize=8.5,
            fontweight="bold", color=CFG["tc_color"])
    ax.text(ann_x + 0.03, tcy - 0.09, r"$T_c$", fontsize=10, color="black", fontstyle="italic")

    # Ts
    ax.plot([tsx + 0.06, ann_x], [tsy - 0.35, tsy - 0.35], ls="--", color=CFG["ts_color"], lw=1.1, zorder=8)
    ax.plot([tsx + 0.06, tsx + 0.06], [tsy, tsy - 0.35], ls="--", color=CFG["ts_color"], lw=1.1, zorder=8)
    ax.text(ann_x + 0.03, tsy - 0.26, "Surface Temperature", fontsize=8.5,
            fontweight="bold", color=CFG["ts_color"])
    ax.text(ann_x + 0.03, tsy - 0.44, r"$T_s$", fontsize=10, color="black", fontstyle="italic")

    # subtitle
    ax.text(cx + 0.15, -1.50,
            "(a) Cylindrical Li-ion Battery\n      Physical Thermal Model",
            fontsize=10.5, fontweight="bold", ha="center", va="top")


# ============================================================
# RIGHT PANEL: Complete thermal equivalent circuit
# ============================================================

def build_right_panel(ax):
    ax.set_xlim(-1.3, 2.0)
    ax.set_ylim(-1.6, 2.0)
    ax.set_aspect("equal")
    ax.axis("off")

    # layout
    lx = -0.55  # left rail x
    rx = 1.55   # right rail x
    ty = 1.25   # top rail y
    by = -0.15  # bottom rail y
    tc_x = 0.10
    ts_x = 0.95

    # rails
    ax.plot([lx, rx], [ty, ty], color=CFG["wire"], lw=2.0, zorder=9)
    ax.plot([lx, rx], [by, by], color=CFG["wire"], lw=2.0, zorder=9)

    # Qe on left
    qm = (ty + by) / 2
    ax.plot([lx, lx], [ty, qm + 0.12], color=CFG["wire"], lw=1.8, zorder=9)
    ax.plot([lx, lx], [qm - 0.12, by], color=CFG["wire"], lw=1.8, zorder=9)
    draw_heat_source(ax, lx, qm, r=0.12, lbl_dx=-0.27)

    # Tc node
    draw_dot(ax, tc_x, ty, CFG["tc_color"], r=0.048, label=r"$T_c$", pos="above")

    # Rc
    draw_resistor_h(ax, tc_x, ts_x, ty, CFG["rc_color"], label=r"$R_c$")

    # Ts node
    draw_dot(ax, ts_x, ty, CFG["ts_color"], r=0.048, label=r"$T_s$", pos="above")

    # Rs
    draw_resistor_h(ax, ts_x, rx, ty, CFG["rs_color"], label=r"$R_s$")

    # right vertical
    ax.plot([rx, rx], [ty, by], color=CFG["wire"], lw=1.8, zorder=9)

    # Cc
    draw_capacitor_v(ax, tc_x, ty - 0.12, by + 0.05, label=r"$C_c$", lbl_side="right")

    # Cs
    draw_capacitor_v(ax, ts_x, ty - 0.12, by + 0.05, label=r"$C_s$", lbl_side="right")

    # ground
    gnd_x = (tc_x + ts_x) / 2
    draw_ground_symbol(ax, gnd_x, by, size=0.10)
    ax.text(gnd_x, by - 0.23, r"$T_a$", fontsize=12, ha="center", va="top")

    # subtitle
    ax.text((lx + rx) / 2, -1.50,
            "(b) Complete Two-Node Thermal\n      Equivalent Circuit",
            fontsize=10.5, fontweight="bold", ha="center", va="top")


# ============================================================
# MAIN
# ============================================================

def main():
    fig, (ax_l, ax_r) = plt.subplots(1, 2, figsize=(CFG["fig_w"], CFG["fig_h"]),
                                      gridspec_kw={"width_ratios": [1, 1.15], "wspace": 0.10})

    build_left_panel(ax_l)
    build_right_panel(ax_r)

    # transition arrow (drawn on right panel extending into left margin)
    ax_r.annotate("", xy=(-0.90, 0.55), xytext=(-1.40, 0.55),
                  arrowprops=dict(arrowstyle="->,head_width=0.35,head_length=0.18",
                                  fc=CFG["arrow_blue"], ec=(0.20, 0.40, 0.80, 0.85), lw=3.5),
                  zorder=20)

    # save
    out_png = "cylindrical_battery_two_node_thermal_model.png"
    fig.savefig(out_png, dpi=CFG["dpi"], bbox_inches="tight", facecolor="white", edgecolor="none")
    print(f"✓ Saved {out_png}")
    plt.close(fig)

    # optional pptx
    try:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        iw, ih = Inches(12.5), Inches(6.0)
        left = (prs.slide_width - iw) // 2
        top = (prs.slide_height - ih) // 2
        slide.shapes.add_picture(out_png, left, top, iw, ih)
        out_pptx = "cylindrical_battery_two_node_thermal_model.pptx"
        prs.save(out_pptx)
        print(f"✓ Saved {out_pptx}")
    except ImportError:
        print("⚠ python-pptx not installed, skipping .pptx generation.")


if __name__ == "__main__":
    main()
