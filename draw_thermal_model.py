"""
Cylindrical 18650 Li-ion Battery Two-Node Thermal Model & Equivalent Circuit
=============================================================================
Generates a publication-quality figure (300 dpi PNG) and an optional PPTX slide.

Usage:
    python draw_thermal_model.py

Output:
    cylindrical_battery_two_node_thermal_model.png
    cylindrical_battery_two_node_thermal_model.pptx  (if python-pptx installed)
"""

import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Arc, Ellipse, Polygon
from matplotlib.lines import Line2D
from matplotlib.collections import PatchCollection
import matplotlib.patheffects as path_effects

# ===========================================================================
# Global style settings
# ===========================================================================
FONT_FAMILY = "serif"
MATH_FONTFAMILY = "cm"  # Computer Modern for LaTeX feel
DPI = 300
FIG_WIDTH = 14  # inches
FIG_HEIGHT = 6.5

# Colors
COLOR_BATTERY_SHELL = (1.0, 0.7, 0.78, 0.55)  # semi-transparent pink
COLOR_BATTERY_CORE = (1.0, 0.85, 0.88, 0.7)
COLOR_TC = (0.5, 0.0, 0.5)  # purple
COLOR_TS = (0.9, 0.1, 0.1)  # red
COLOR_RC = (1.0, 0.75, 0.8)  # light pink
COLOR_RS = (1.0, 0.6, 0.2)  # orange
COLOR_WIRE = "black"
COLOR_ARROW = (0.3, 0.6, 0.9, 0.7)  # light blue
COLOR_GROUND = "black"
COLOR_METAL = (0.75, 0.75, 0.78)
COLOR_METAL_HIGHLIGHT = (0.92, 0.92, 0.95)
COLOR_INSULATION = (0.15, 0.15, 0.15)

plt.rcParams.update({
    "font.family": FONT_FAMILY,
    "mathtext.fontset": MATH_FONTFAMILY,
    "font.size": 10,
    "axes.linewidth": 0.0,
})


# ===========================================================================
# Helper drawing functions
# ===========================================================================

def draw_cylinder_body(ax, cx, cy, width, height, tilt=0.18):
    """Draw a 3D-looking cylinder (battery body) with transparency."""
    # Ellipse semi-axes
    a = width / 2
    b = a * tilt

    # Bottom ellipse
    bottom_ellipse = Ellipse((cx, cy - height / 2), width, 2 * b,
                             facecolor=COLOR_BATTERY_SHELL, edgecolor=(0.6, 0.3, 0.4, 0.8),
                             linewidth=1.2, zorder=2)
    ax.add_patch(bottom_ellipse)

    # Cylinder side (rectangle with gradient-like effect using multiple rects)
    n_strips = 40
    for i in range(n_strips):
        frac = i / n_strips
        x_left = cx - a + frac * width
        strip_w = width / n_strips
        # Simulate highlight in center
        dist_from_center = abs(frac - 0.5) / 0.5
        alpha = 0.35 + 0.25 * (1 - dist_from_center ** 2)
        color = (1.0, 0.55 + 0.2 * (1 - dist_from_center), 0.7 + 0.1 * (1 - dist_from_center), alpha)
        rect = plt.Rectangle((x_left, cy - height / 2), strip_w, height,
                              facecolor=color, edgecolor="none", zorder=1)
        ax.add_patch(rect)

    # Side outline
    ax.plot([cx - a, cx - a], [cy - height / 2, cy + height / 2], color=(0.6, 0.3, 0.4, 0.8), lw=1.2, zorder=3)
    ax.plot([cx + a, cx + a], [cy - height / 2, cy + height / 2], color=(0.6, 0.3, 0.4, 0.8), lw=1.2, zorder=3)

    # Top ellipse (main face)
    top_ellipse = Ellipse((cx, cy + height / 2), width, 2 * b,
                          facecolor=(1.0, 0.82, 0.88, 0.85), edgecolor=(0.6, 0.3, 0.4, 0.8),
                          linewidth=1.2, zorder=4)
    ax.add_patch(top_ellipse)

    # Core region (inner lighter cylinder hint)
    core_w = width * 0.6
    core_h = height * 0.85
    core_a = core_w / 2
    for i in range(20):
        frac = i / 20
        x_left = cx - core_a + frac * core_w
        strip_w = core_w / 20
        dist_from_center = abs(frac - 0.5) / 0.5
        alpha = 0.15 + 0.15 * (1 - dist_from_center ** 2)
        color = (0.95, 0.8, 0.85, alpha)
        rect = plt.Rectangle((x_left, cy - core_h / 2), strip_w, core_h,
                              facecolor=color, edgecolor="none", zorder=1.5)
        ax.add_patch(rect)

    return cx, cy, a, b, height


def draw_positive_terminal(ax, cx, cy_top, a, b):
    """
    Draw the positive terminal cap on top of the battery.
    
    Improved version: the cap sits directly on top of the battery cylinder
    with a visible 3D side wall (short cylinder), metallic gradient shading,
    insulation ring, and a raised central bump with specular highlight.
    """
    # --- Parameters ---
    cap_side_height = 0.10  # height of the cap side wall
    cap_radius_factor = 0.92  # slightly smaller than battery top
    cap_a = a * cap_radius_factor  # half-width of cap ellipse
    cap_b = b * cap_radius_factor  # half-height (perspective)

    # The cap bottom sits exactly at cy_top (battery top ellipse center)
    cap_bottom_y = cy_top
    cap_top_y = cap_bottom_y + cap_side_height

    # --- 1. Side wall of the cap (short metallic cylinder) ---
    # Draw gradient strips for metallic look
    n_strips = 30
    for i in range(n_strips):
        frac = i / n_strips
        x_left = cx - cap_a + frac * (2 * cap_a)
        strip_w = (2 * cap_a) / n_strips
        # Metallic shading: brighter in center, darker at edges
        dist_from_center = abs(frac - 0.5) / 0.5
        brightness = 0.82 - 0.25 * dist_from_center ** 1.5
        color = (brightness, brightness, brightness + 0.03, 0.95)
        rect = plt.Rectangle((x_left, cap_bottom_y), strip_w, cap_side_height,
                              facecolor=color, edgecolor="none", zorder=4.5)
        ax.add_patch(rect)

    # Side wall outline (left and right edges)
    ax.plot([cx - cap_a, cx - cap_a], [cap_bottom_y, cap_top_y],
            color=(0.35, 0.35, 0.38), lw=1.3, zorder=4.6)
    ax.plot([cx + cap_a, cx + cap_a], [cap_bottom_y, cap_top_y],
            color=(0.35, 0.35, 0.38), lw=1.3, zorder=4.6)

    # --- 2. Top face of the cap (outer metal ring) ---
    outer_ring = Ellipse((cx, cap_top_y), 2 * cap_a, 2 * cap_b,
                         facecolor=(0.78, 0.78, 0.81), edgecolor=(0.35, 0.35, 0.38),
                         linewidth=1.2, zorder=5)
    ax.add_patch(outer_ring)

    # Specular highlight on outer ring (upper-left)
    hl_offset_x = -cap_a * 0.25
    hl_offset_y = cap_b * 0.2
    highlight = Ellipse((cx + hl_offset_x, cap_top_y + hl_offset_y),
                        cap_a * 0.7, cap_b * 0.55,
                        facecolor=(0.95, 0.95, 0.98), edgecolor="none",
                        alpha=0.5, zorder=5.05)
    ax.add_patch(highlight)

    # --- 3. Black insulation ring ---
    insul_a = cap_a * 0.72
    insul_b = cap_b * 0.72
    insul_ring = Ellipse((cx, cap_top_y), 2 * insul_a, 2 * insul_b,
                         facecolor=(0.12, 0.12, 0.12), edgecolor=(0.05, 0.05, 0.05),
                         linewidth=0.9, zorder=5.1)
    ax.add_patch(insul_ring)

    # --- 4. Central raised positive bump ---
    bump_a = cap_a * 0.38
    bump_b = cap_b * 0.50
    # Bump base (slightly darker ring to suggest elevation)
    bump_base = Ellipse((cx, cap_top_y + 0.01), 2 * bump_a * 1.1, 2 * bump_b * 1.1,
                        facecolor=(0.6, 0.6, 0.63), edgecolor=(0.4, 0.4, 0.42),
                        linewidth=0.7, zorder=5.2)
    ax.add_patch(bump_base)

    # Bump top surface
    bump_top = Ellipse((cx, cap_top_y + 0.025), 2 * bump_a, 2 * bump_b,
                       facecolor=(0.82, 0.82, 0.86), edgecolor=(0.5, 0.5, 0.53),
                       linewidth=0.8, zorder=5.3)
    ax.add_patch(bump_top)

    # Bump specular highlight
    bump_hl = Ellipse((cx - bump_a * 0.2, cap_top_y + 0.035),
                      bump_a * 0.6, bump_b * 0.5,
                      facecolor=(0.98, 0.98, 1.0), edgecolor="none",
                      alpha=0.7, zorder=5.4)
    ax.add_patch(bump_hl)

    # --- 5. "+" label above the terminal ---
    ax.text(cx, cap_top_y + cap_b + 0.06, "+", fontsize=10, fontweight="bold",
            ha="center", va="center", color="black", zorder=6)


def draw_heat_source_symbol(ax, x, y, radius=0.08, label=r"$Q_e$", label_offset=(-0.18, 0)):
    """Draw a heat source symbol (circle with sine wave inside)."""
    circle = plt.Circle((x, y), radius, facecolor="white", edgecolor="black", linewidth=1.5, zorder=10)
    ax.add_patch(circle)
    # Sine wave inside
    t = np.linspace(-0.7, 0.7, 30) * radius
    wave = np.sin(t / radius * 2.5 * np.pi) * radius * 0.35
    ax.plot(t + x, wave + y, color="black", lw=1.2, zorder=11)
    # Label
    ax.text(x + label_offset[0], y + label_offset[1], label, fontsize=11, ha="center", va="center",
            color="black", zorder=11)


def draw_resistor(ax, x1, y1, x2, y2, color=COLOR_RC, label="", label_side="above"):
    """Draw a resistor (rectangle) between two points along a line."""
    mx, my = (x1 + x2) / 2, (y1 + y2) / 2
    dx, dy = x2 - x1, y2 - y1
    length = np.sqrt(dx ** 2 + dy ** 2)
    angle = np.degrees(np.arctan2(dy, dx))

    res_len = length * 0.5
    res_h = 0.09

    # Draw connecting wires
    ax.plot([x1, x1 + dx * 0.25], [y1, y1 + dy * 0.25], color=COLOR_WIRE, lw=1.8, zorder=9)
    ax.plot([x1 + dx * 0.75, x2], [y1 + dy * 0.75, y2], color=COLOR_WIRE, lw=1.8, zorder=9)

    # Resistor box
    from matplotlib.transforms import Affine2D
    rect = FancyBboxPatch((mx - res_len / 2, my - res_h / 2), res_len, res_h,
                          boxstyle="round,pad=0.01", facecolor=color, edgecolor="black",
                          linewidth=1.3, zorder=10)
    # Rotate if needed
    if abs(angle) > 1 and abs(angle - 0) > 1:
        t = Affine2D().rotate_deg_around(mx, my, angle) + ax.transData
        rect.set_transform(t)
    ax.add_patch(rect)

    # Label
    if label:
        offset = 0.12 if label_side == "above" else -0.12
        if abs(angle) < 5:  # horizontal
            ax.text(mx, my + offset, label, fontsize=11, ha="center", va="center", color="black", zorder=11)
        else:  # vertical
            ax.text(mx + offset, my, label, fontsize=11, ha="center", va="center", color="black", zorder=11)


def draw_capacitor(ax, x, y_top, y_bot, label="", label_side="right"):
    """Draw a capacitor symbol (two parallel lines) vertically."""
    gap = 0.06
    plate_w = 0.12
    mid_y = (y_top + y_bot) / 2

    # Wires
    ax.plot([x, x], [y_top, mid_y + gap], color=COLOR_WIRE, lw=1.8, zorder=9)
    ax.plot([x, x], [mid_y - gap, y_bot], color=COLOR_WIRE, lw=1.8, zorder=9)

    # Plates
    ax.plot([x - plate_w, x + plate_w], [mid_y + gap, mid_y + gap], color="black", lw=2.5, zorder=10)
    ax.plot([x - plate_w, x + plate_w], [mid_y - gap, mid_y - gap], color="black", lw=2.5, zorder=10)

    # Label
    if label:
        offset = 0.16 if label_side == "right" else -0.16
        ax.text(x + offset, mid_y, label, fontsize=11, ha="left" if label_side == "right" else "right",
                va="center", color="black", zorder=11)


def draw_ground(ax, x, y, size=0.08):
    """Draw a ground symbol."""
    ax.plot([x, x], [y, y - size * 0.5], color=COLOR_WIRE, lw=1.8, zorder=9)
    for i, frac in enumerate([1.0, 0.65, 0.3]):
        w = size * frac
        yy = y - size * 0.5 - i * size * 0.25
        ax.plot([x - w, x + w], [yy, yy], color=COLOR_WIRE, lw=1.8, zorder=9)


def draw_node_dot(ax, x, y, color="black", radius=0.04, label="", label_pos="above"):
    """Draw a node dot with optional label."""
    dot = plt.Circle((x, y), radius, facecolor=color, edgecolor="black", linewidth=0.8, zorder=12)
    ax.add_patch(dot)
    if label:
        if label_pos == "above":
            ax.text(x, y + 0.13, label, fontsize=11, ha="center", va="center", color="black", zorder=12)
        elif label_pos == "below":
            ax.text(x, y - 0.13, label, fontsize=11, ha="center", va="center", color="black", zorder=12)
        elif label_pos == "left":
            ax.text(x - 0.13, y, label, fontsize=11, ha="right", va="center", color="black", zorder=12)
        elif label_pos == "right":
            ax.text(x + 0.13, y, label, fontsize=11, ha="left", va="center", color="black", zorder=12)


# ===========================================================================
# Main figure composition
# ===========================================================================

def create_figure():
    fig, (ax_left, ax_right) = plt.subplots(1, 2, figsize=(FIG_WIDTH, FIG_HEIGHT),
                                             gridspec_kw={"width_ratios": [1, 1.15], "wspace": 0.12})
    for ax in (ax_left, ax_right):
        ax.set_xlim(-1.2, 1.8)
        ax.set_ylim(-1.3, 1.8)
        ax.set_aspect("equal")
        ax.axis("off")

    # ===================================================================
    # LEFT: Battery physical model
    # ===================================================================
    bx, by = 0.0, 0.2  # battery center
    bw, bh = 1.0, 2.2  # width, height
    tilt = 0.18

    cx, cy, a, b, height = draw_cylinder_body(ax_left, bx, by, bw, bh, tilt)
    draw_positive_terminal(ax_left, bx, by + bh / 2, a, b * 0.9)

    # --- Internal thermal network (overlaid on battery) ---
    # Positions for internal nodes
    qe_x, qe_y = bx - 0.15, by - 0.55  # heat source
    tc_x, tc_y = bx - 0.15, by + 0.25   # Tc node
    ts_x, ts_y = bx + 0.30, by + 0.25   # Ts node
    gnd_y = by - 0.95

    # Qe symbol
    draw_heat_source_symbol(ax_left, qe_x, qe_y, radius=0.09,
                            label=r"$Q_e$", label_offset=(-0.22, 0))

    # Wire Qe to Tc
    ax_left.plot([qe_x, qe_x], [qe_y + 0.09, tc_y], color=COLOR_WIRE, lw=1.5, zorder=9)

    # Tc node
    draw_node_dot(ax_left, tc_x, tc_y, color=COLOR_TC, radius=0.04,
                  label=r"$T_c$", label_pos="above")

    # Ts node
    draw_node_dot(ax_left, ts_x, ts_y, color=COLOR_TS, radius=0.04,
                  label=r"$T_s$", label_pos="above")

    # Rc between Tc and Ts
    draw_resistor(ax_left, tc_x, tc_y, ts_x, ts_y, color=COLOR_RC, label=r"$R_c$", label_side="above")

    # Cc from Tc downward
    cc_bot = gnd_y
    draw_capacitor(ax_left, tc_x, tc_y - 0.08, cc_bot + 0.05, label=r"$C_c$", label_side="left")

    # Rs from Ts downward
    rs_bot = gnd_y + 0.35
    draw_resistor(ax_left, ts_x, ts_y, ts_x, rs_bot, color=COLOR_RS, label=r"$R_s$", label_side="right")

    # Cs from Ts downward (parallel path)
    cs_x = ts_x + 0.25
    ax_left.plot([ts_x, cs_x], [ts_y, ts_y], color=COLOR_WIRE, lw=1.2, zorder=9)
    draw_capacitor(ax_left, cs_x, ts_y - 0.08, cc_bot + 0.05, label=r"$C_s$", label_side="right")
    ax_left.plot([cs_x, cs_x], [cc_bot + 0.05, cc_bot], color=COLOR_WIRE, lw=1.2, zorder=9)

    # Wire Rs bottom to ground line
    ax_left.plot([ts_x, ts_x], [rs_bot, cc_bot], color=COLOR_WIRE, lw=1.2, zorder=9)

    # Ground line
    ax_left.plot([tc_x - 0.1, cs_x + 0.1], [cc_bot, cc_bot], color=COLOR_WIRE, lw=1.5, zorder=9)

    # Qe bottom to ground
    ax_left.plot([qe_x, qe_x], [qe_y - 0.09, cc_bot], color=COLOR_WIRE, lw=1.5, zorder=9)

    # Ground symbol and Ta
    gnd_x = (tc_x + ts_x) / 2
    draw_ground(ax_left, gnd_x, cc_bot, size=0.09)
    ax_left.text(gnd_x, cc_bot - 0.18, r"$T_a$", fontsize=11, ha="center", va="top", color="black")

    # --- Annotation dashed lines ---
    ann_x = 1.05
    # Tc annotation
    ax_left.annotate("", xy=(ann_x, tc_y), xytext=(tc_x + 0.06, tc_y),
                     arrowprops=dict(arrowstyle="-", linestyle="dashed", color=COLOR_TC, lw=1.2))
    ax_left.text(ann_x + 0.02, tc_y + 0.08, "Internal Temperature", fontsize=8.5,
                 fontweight="bold", color=COLOR_TC, va="center")
    ax_left.text(ann_x + 0.02, tc_y - 0.08, r"$T_c$", fontsize=10, fontstyle="italic",
                 color="black", va="center")

    # Ts annotation
    ax_left.annotate("", xy=(ann_x, ts_y - 0.35), xytext=(ts_x + 0.06, ts_y - 0.05),
                     arrowprops=dict(arrowstyle="-", linestyle="dashed", color=COLOR_TS, lw=1.2))
    ax_left.text(ann_x + 0.02, ts_y - 0.27, "Surface Temperature", fontsize=8.5,
                 fontweight="bold", color=COLOR_TS, va="center")
    ax_left.text(ann_x + 0.02, ts_y - 0.43, r"$T_s$", fontsize=10, fontstyle="italic",
                 color="black", va="center")

    # Subtitle (a)
    ax_left.text(bx + 0.1, -1.25, "(a) Cylindrical Li-ion Battery\n     Physical Thermal Model",
                 fontsize=10, fontweight="bold", ha="center", va="top", color="black")

    # ===================================================================
    # RIGHT: Complete two-node thermal equivalent circuit
    # ===================================================================
    # Layout coordinates
    left_x = -0.6
    tc_cx = 0.0
    ts_cx = 0.8
    right_x = 1.4
    top_y = 1.2
    bot_y = -0.2
    gnd_cy = bot_y

    # Top rail
    ax_right.plot([left_x, right_x], [top_y, top_y], color=COLOR_WIRE, lw=2.0, zorder=9)
    # Bottom rail
    ax_right.plot([left_x, right_x], [bot_y, bot_y], color=COLOR_WIRE, lw=2.0, zorder=9)

    # Heat source Qe on left
    qe_mid_y = (top_y + bot_y) / 2
    ax_right.plot([left_x, left_x], [top_y, qe_mid_y + 0.12], color=COLOR_WIRE, lw=1.8, zorder=9)
    ax_right.plot([left_x, left_x], [qe_mid_y - 0.12, bot_y], color=COLOR_WIRE, lw=1.8, zorder=9)
    draw_heat_source_symbol(ax_right, left_x, qe_mid_y, radius=0.11,
                            label=r"$Q_e$", label_offset=(-0.25, 0))

    # Tc node on top rail
    draw_node_dot(ax_right, tc_cx, top_y, color=COLOR_TC, radius=0.045,
                  label=r"$T_c$", label_pos="above")

    # Rc between Tc and Ts on top rail
    draw_resistor(ax_right, tc_cx, top_y, ts_cx, top_y, color=COLOR_RC, label=r"$R_c$", label_side="above")

    # Ts node on top rail
    draw_node_dot(ax_right, ts_cx, top_y, color=COLOR_TS, radius=0.045,
                  label=r"$T_s$", label_pos="above")

    # Rs between Ts and right end on top rail
    draw_resistor(ax_right, ts_cx, top_y, right_x, top_y, color=COLOR_RS, label=r"$R_s$", label_side="above")

    # Right vertical wire
    ax_right.plot([right_x, right_x], [top_y, bot_y], color=COLOR_WIRE, lw=1.8, zorder=9)

    # Cc from Tc to bottom
    draw_capacitor(ax_right, tc_cx, top_y - 0.1, bot_y + 0.05, label=r"$C_c$", label_side="right")

    # Cs from Ts to bottom
    draw_capacitor(ax_right, ts_cx, top_y - 0.1, bot_y + 0.05, label=r"$C_s$", label_side="right")

    # Ground symbol at center of bottom rail
    gnd_cx = (tc_cx + ts_cx) / 2
    draw_ground(ax_right, gnd_cx, bot_y, size=0.10)
    ax_right.text(gnd_cx, bot_y - 0.22, r"$T_a$", fontsize=12, ha="center", va="top", color="black")

    # Subtitle (b)
    ax_right.text((left_x + right_x) / 2, -1.25,
                  "(b) Complete Two-Node Thermal\n     Equivalent Circuit",
                  fontsize=10, fontweight="bold", ha="center", va="top", color="black")

    # ===================================================================
    # CENTER: Transition arrow between left and right
    # ===================================================================
    ax_right.annotate(
        "", xy=(-0.95, 0.5), xytext=(-1.45, 0.5),
        arrowprops=dict(arrowstyle="->,head_width=0.4,head_length=0.2",
                        fc=COLOR_ARROW, ec=(0.2, 0.4, 0.8, 0.8), lw=3),
        zorder=20
    )

    return fig


def save_figure(fig, filename="cylindrical_battery_two_node_thermal_model.png"):
    """Save figure as high-resolution PNG."""
    fig.savefig(filename, dpi=DPI, bbox_inches="tight", facecolor="white", edgecolor="none")
    print(f"Saved: {filename}")


def create_pptx(image_path="cylindrical_battery_two_node_thermal_model.png",
                output_path="cylindrical_battery_two_node_thermal_model.pptx"):
    """Create a 16:9 PPTX slide with the figure inserted."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
    except ImportError:
        print("python-pptx not installed. Skipping PPTX generation.")
        print("Install with: pip install python-pptx")
        return

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Center the image on slide
    img_width = Inches(12.0)
    img_height = Inches(5.6)
    left = (prs.slide_width - img_width) // 2
    top = (prs.slide_height - img_height) // 2
    slide.shapes.add_picture(image_path, left, top, img_width, img_height)

    prs.save(output_path)
    print(f"Saved: {output_path}")


# ===========================================================================
# Entry point
# ===========================================================================

if __name__ == "__main__":
    fig = create_figure()
    save_figure(fig)
    plt.close(fig)
    create_pptx()
